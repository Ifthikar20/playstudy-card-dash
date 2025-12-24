"""
Study Sessions API endpoints with AI content processing.
"""
from fastapi import APIRouter, Depends, HTTPException, Request
from pydantic import BaseModel, Field
from typing import List, Optional
from sqlalchemy.orm import Session
from datetime import datetime
from openai import OpenAI
from anthropic import Anthropic
import json
import io
import uuid
import logging
import base64
import tempfile
import subprocess
import os
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation

from app.config import settings
from app.dependencies import get_current_active_user, get_db
from app.models.user import User
from app.models.study_session import StudySession
from app.models.topic import Topic
from app.models.question import Question
from app.core.rate_limit import limiter

logger = logging.getLogger(__name__)

router = APIRouter()


def build_topic_hierarchy(session: StudySession, db: Session) -> List[TopicSchema]:
    """
    Build hierarchical topic structure with questions for a study session.
    OPTIMIZED: Uses pre-loaded relationships when available to avoid N+1 queries.

    Args:
        session: StudySession database model (may have eager-loaded topics)
        db: Database session

    Returns:
        List of TopicSchema objects with nested subtopics and questions
    """
    from app.models.question import Question

    # OPTIMIZED: Check if topics are already loaded (from eager loading)
    # If session.topics is already populated, use it instead of querying
    if hasattr(session, 'topics') and session.topics:
        all_topics = sorted(session.topics, key=lambda t: t.order_index or 0)
    else:
        # Fallback: Query with eager loading for questions
        from sqlalchemy.orm import selectinload
        all_topics = db.query(Topic).options(
            selectinload(Topic.questions)
        ).filter(
            Topic.study_session_id == session.id
        ).order_by(Topic.order_index).all()

    # Build hierarchical structure
    categories = [t for t in all_topics if t.is_category and t.parent_topic_id is None]

    result_topics = []

    for cat_idx, category in enumerate(categories):
        # Get subtopics for this category
        subtopics = [t for t in all_topics if t.parent_topic_id == category.id]

        category_schema = TopicSchema(
            id=f"category-{cat_idx+1}",
            db_id=category.id,
            title=category.title,
            description=category.description or "",
            isCategory=True,
            parentTopicId=None,
            questions=[],
            subtopics=[]
        )

        for sub_idx, subtopic in enumerate(subtopics):
            # OPTIMIZED: Use pre-loaded questions if available
            if hasattr(subtopic, 'questions') and subtopic.questions:
                questions = sorted(subtopic.questions, key=lambda q: q.order_index or 0)
            else:
                # Fallback: Query if not pre-loaded
                questions = db.query(Question).filter(
                    Question.topic_id == subtopic.id
                ).order_by(Question.order_index).all()

            questions_list = [
                QuestionSchema(
                    id=f"topic-{sub_idx+1}-q{q.order_index+1}",
                    question=q.question,
                    options=q.options,
                    correctAnswer=q.correct_answer,
                    explanation=q.explanation,
                    sourceText=q.source_text,
                    sourcePage=q.source_page
                )
                for q in questions
            ]

            subtopic_schema = TopicSchema(
                id=f"subtopic-{cat_idx+1}-{sub_idx+1}",
                db_id=subtopic.id,
                title=subtopic.title,
                description=subtopic.description or "",
                questions=questions_list,
                completed=subtopic.completed or False,
                score=subtopic.score,
                currentQuestionIndex=subtopic.current_question_index or 0,
                isCategory=False,
                parentTopicId=f"category-{cat_idx+1}",
                subtopics=[]
            )
            category_schema.subtopics.append(subtopic_schema)

        result_topics.append(category_schema)

    return result_topics


def analyze_content_complexity(text: str) -> dict:
    """
    Analyze content complexity and recommend topic/question counts.

    Args:
        text: Extracted text content

    Returns:
        Dictionary with analysis results including:
        - word_count: Total number of words
        - estimated_reading_time: Time to read in minutes
        - recommended_topics: Suggested number of topics
        - recommended_questions: Suggested questions per topic
        - complexity_score: 0-1 score of content complexity
        - unique_word_ratio: Vocabulary richness
    """
    words = text.split()
    word_count = len(words)

    # Calculate unique words ratio (vocabulary richness)
    unique_words = len(set(word.lower() for word in words if word.isalnum()))
    unique_word_ratio = unique_words / max(word_count, 1)

    # Average word length (longer words = more complex)
    avg_word_length = sum(len(word) for word in words) / max(word_count, 1)

    # Sentence count (approximate by counting periods, exclamation, question marks)
    sentences = len([c for c in text if c in '.!?'])
    avg_sentence_length = word_count / max(sentences, 1)

    # Calculate complexity score (0-1 scale)
    # Based on: vocabulary richness, average word length, sentence complexity
    complexity_score = min(1.0, (
        (unique_word_ratio * 0.4) +  # 40% weight on vocabulary
        (min(avg_word_length / 8, 1.0) * 0.3) +  # 30% weight on word length
        (min(avg_sentence_length / 25, 1.0) * 0.3)  # 30% weight on sentence length
    ))

    # Reading time (average reading speed: 200-250 words/minute)
    estimated_reading_time = max(1, round(word_count / 225))

    # Recommend topics based on word count and complexity
    # Focus on quality over quantity - only create topics that will have sufficient questions
    # Very short (< 100 words): 1 topic
    # Short (100-500 words): 2-3 topics
    # Medium (500-2000 words): 3-6 topics
    # Long (2000-5000 words): 6-10 topics
    # Very long (5000-10000 words): 10-15 topics
    # Extremely long (10000-20000 words): 15-25 topics
    # Massive (20000+ words): 25-35 topics (maximum to ensure quality)
    if word_count < 100:
        base_topics = 1
    elif word_count < 500:
        base_topics = 2
    elif word_count < 2000:
        base_topics = 4
    elif word_count < 5000:
        base_topics = 8
    elif word_count < 10000:
        base_topics = 12
    elif word_count < 20000:
        base_topics = 20
    else:
        base_topics = 30

    # Adjust based on complexity (more conservative multiplier)
    recommended_topics = max(1, min(35, round(base_topics * (0.8 + complexity_score * 0.4))))

    # Recommend questions per topic based on content depth
    # Aim for maximum questions to ensure comprehensive coverage
    # More questions = better learning coverage and understanding
    if word_count < 1000:
        base_questions = 15
    elif word_count < 3000:
        base_questions = 20
    elif word_count < 10000:
        base_questions = 25
    else:
        base_questions = 30

    # Adjust based on complexity (allow up to 100 questions for complex topics)
    recommended_questions = max(10, min(100, round(base_questions * (0.9 + complexity_score * 0.6))))

    return {
        'word_count': word_count,
        'estimated_reading_time': estimated_reading_time,
        'recommended_topics': recommended_topics,
        'recommended_questions': recommended_questions,
        'complexity_score': round(complexity_score, 2),
        'unique_word_ratio': round(unique_word_ratio, 2),
        'avg_word_length': round(avg_word_length, 1),
        'avg_sentence_length': round(avg_sentence_length, 1)
    }


def detect_file_type_and_extract(content: str) -> tuple[str, str, str]:
    """
    Detect file type and extract text from uploaded content.

    Returns:
        tuple: (extracted_text, file_type, original_base64_content)
            - extracted_text: The text content extracted from the file
            - file_type: One of: 'pdf', 'pptx', 'docx', 'txt'
            - original_base64_content: The original base64 encoded file
    """
    import base64

    content_bytes = None
    file_type = 'txt'
    original_base64 = content

    # First, aggressively try base64 decode (most likely for file uploads)
    try:
        # Try base64 decode - this is the most common format for file uploads via JSON
        decoded = base64.b64decode(content, validate=False)
        # Check if it looks like a valid file (ZIP/docx or PDF)
        if decoded.startswith(b'PK\x03\x04') or decoded.startswith(b'%PDF'):
            content_bytes = decoded
    except Exception:
        pass

    # Process based on file type if we have bytes
    if content_bytes:
        # Office documents (ZIP/docx/pptx)
        if content_bytes.startswith(b'PK\x03\x04'):
            # Try PowerPoint first
            try:
                prs = Presentation(io.BytesIO(content_bytes))
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_parts.append(shape.text)
                text = '\n'.join(text_parts)
                if text.strip():
                    return (text, 'pptx', original_base64)
                raise ValueError("No text content found in PowerPoint")
            except Exception as pptx_error:
                # If PowerPoint fails, try Word document
                try:
                    doc = Document(io.BytesIO(content_bytes))
                    text = '\n'.join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])
                    if text.strip():
                        return (text, 'docx', original_base64)
                    raise ValueError("No text content found in Word document")
                except Exception as docx_error:
                    raise HTTPException(
                        status_code=400,
                        detail=f"Failed to extract text from Office document. PowerPoint error: {str(pptx_error)}. Word error: {str(docx_error)}"
                    )

        # PDF file
        elif content_bytes.startswith(b'%PDF'):
            try:
                pdf_reader = PdfReader(io.BytesIO(content_bytes))
                text = '\n'.join([page.extract_text() for page in pdf_reader.pages])
                if text.strip():
                    return (text, 'pdf', original_base64)
                raise ValueError("No text content found in PDF")
            except Exception as e:
                raise HTTPException(
                    status_code=400,
                    detail=f"Failed to extract text from PDF: {str(e)}"
                )

    # If we got here, treat as plain text
    cleaned_text = content.replace('\ufffd', '').strip()
    if cleaned_text and len(cleaned_text) > 10:
        return (cleaned_text, 'txt', original_base64)

    # Last resort - return original
    return (content, 'txt', original_base64)


def extract_text_from_content(content: str) -> str:
    """
    Extract readable text from uploaded content.
    Handles:
    - Plain text
    - Word documents (.docx)
    - PowerPoint presentations (.pptx, .ppt)
    - PDF files
    - Base64 encoded content

    Args:
        content: Raw content string (may be binary, text, or base64)

    Returns:
        Extracted text string
    """
    import base64

    content_bytes = None

    # First, aggressively try base64 decode (most likely for file uploads)
    try:
        # Try base64 decode - this is the most common format for file uploads via JSON
        decoded = base64.b64decode(content, validate=False)
        # Check if it looks like a valid file (ZIP/docx or PDF)
        if decoded.startswith(b'PK\x03\x04') or decoded.startswith(b'%PDF'):
            content_bytes = decoded
    except Exception:
        pass

    # If base64 didn't work and content looks like binary, try encoding it
    if content_bytes is None and (content.startswith('PK\x03\x04') or content.startswith('%PDF')):
        # Try different encodings to convert string to bytes
        for encoding in ['latin-1', 'iso-8859-1', 'cp1252', 'utf-8']:
            try:
                content_bytes = content.encode(encoding)
                break
            except (UnicodeEncodeError, UnicodeDecodeError):
                continue

    # Process based on file type if we have bytes
    if content_bytes:
        # Office documents (ZIP/docx/pptx)
        if content_bytes.startswith(b'PK\x03\x04'):
            # Try PowerPoint first
            try:
                prs = Presentation(io.BytesIO(content_bytes))
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_parts.append(shape.text)
                text = '\n'.join(text_parts)
                if text.strip():
                    return text
                raise ValueError("No text content found in PowerPoint")
            except Exception as pptx_error:
                # If PowerPoint fails, try Word document
                try:
                    doc = Document(io.BytesIO(content_bytes))
                    text = '\n'.join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])
                    if text.strip():
                        return text
                    raise ValueError("No text content found in Word document")
                except Exception as docx_error:
                    raise HTTPException(
                        status_code=400,
                        detail=f"Failed to extract text from Office document. PowerPoint error: {str(pptx_error)}. Word error: {str(docx_error)}"
                    )

        # PDF file
        elif content_bytes.startswith(b'%PDF'):
            try:
                pdf_reader = PdfReader(io.BytesIO(content_bytes))
                text = '\n'.join([page.extract_text() for page in pdf_reader.pages])
                if text.strip():
                    return text
                raise ValueError("No text content found in PDF")
            except Exception as e:
                raise HTTPException(
                    status_code=400,
                    detail=f"Failed to extract text from PDF: {str(e)}"
                )

    # If we got here, treat as plain text
    # Clean up any Unicode replacement characters
    cleaned_text = content.replace('\ufffd', '').strip()
    if cleaned_text and len(cleaned_text) > 10:
        return cleaned_text

    # Last resort - return original
    return content


def convert_pptx_to_pdf(pptx_base64: str) -> Optional[str]:
    """
    Convert PowerPoint (PPTX) file to PDF using LibreOffice.

    Args:
        pptx_base64: Base64-encoded PPTX file content

    Returns:
        Base64-encoded PDF file content, or None if conversion fails
    """
    temp_dir = None
    try:
        # Create temporary directory for conversion
        temp_dir = tempfile.mkdtemp()
        pptx_path = os.path.join(temp_dir, "presentation.pptx")
        pdf_path = os.path.join(temp_dir, "presentation.pdf")

        # Decode and write PPTX file
        pptx_bytes = base64.b64decode(pptx_base64)
        with open(pptx_path, 'wb') as f:
            f.write(pptx_bytes)

        # Convert using LibreOffice
        # --headless: Run without GUI
        # --convert-to pdf: Convert to PDF format
        # --outdir: Output directory
        try:
            result = subprocess.run(
                ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', temp_dir, pptx_path],
                capture_output=True,
                text=True,
                timeout=30  # 30 second timeout
            )

            # Check if conversion succeeded
            if result.returncode != 0:
                logger.warning(f"LibreOffice conversion failed: {result.stderr}")
                return None

            # Check if PDF file was created
            if not os.path.exists(pdf_path):
                logger.warning("LibreOffice conversion completed but PDF file not found")
                return None

            # Read and encode PDF
            with open(pdf_path, 'rb') as f:
                pdf_bytes = f.read()

            pdf_base64 = base64.b64encode(pdf_bytes).decode('utf-8')
            logger.info(f"Successfully converted PPTX to PDF ({len(pdf_bytes)} bytes)")
            return pdf_base64

        except subprocess.TimeoutExpired:
            logger.warning("LibreOffice conversion timed out after 30 seconds")
            return None
        except FileNotFoundError:
            logger.warning("LibreOffice not found - PPTX to PDF conversion unavailable")
            return None

    except Exception as e:
        logger.error(f"Error converting PPTX to PDF: {e}")
        return None
    finally:
        # Clean up temporary files
        if temp_dir and os.path.exists(temp_dir):
            try:
                import shutil
                shutil.rmtree(temp_dir)
            except Exception as e:
                logger.warning(f"Failed to clean up temp directory: {e}")


class QuestionSchema(BaseModel):
    """Schema for a quiz question."""
    id: str
    question: str
    options: List[str]
    correctAnswer: int
    explanation: str
    sourceText: Optional[str] = None  # Source text snippet from document
    sourcePage: Optional[int] = None  # Page number in source document


class TopicSchema(BaseModel):
    """Schema for a study topic."""
    id: str
    db_id: Optional[int] = None  # Database ID for syncing progress
    title: str
    description: str
    questions: List[QuestionSchema] = []
    completed: bool = False
    score: Optional[int] = None
    currentQuestionIndex: int = 0
    isCategory: bool = False
    parentTopicId: Optional[str] = None
    subtopics: List['TopicSchema'] = []

# Enable forward references for recursive schema
TopicSchema.model_rebuild()


class CreateStudySessionRequest(BaseModel):
    """Request schema for creating a study session."""
    title: str = Field(..., min_length=1, max_length=200)
    content: str = Field(..., min_length=10, max_length=100000000)  # 100MB limit for base64 encoded files (large PDFs)
    num_topics: int = Field(default=4, ge=1, le=100)  # Dynamic: 1-100 topics based on content size
    questions_per_topic: int = Field(default=50, ge=5, le=200)  # Lowered minimum to 5 to support flexible question generation
    progressive_load: bool = Field(default=False)  # DISABLED: Generate ALL questions upfront for better UX


class AnalyzeContentRequest(BaseModel):
    """Request schema for analyzing content before creating a session."""
    content: str = Field(..., min_length=10, max_length=100000000)  # 100MB limit


class ContentAnalysisResponse(BaseModel):
    """Response schema for content analysis."""
    word_count: int
    estimated_reading_time: int  # in minutes
    recommended_topics: int
    recommended_questions: int
    complexity_score: float  # 0-1 scale
    content_summary: str


class CreateStudySessionResponse(BaseModel):
    """Response schema for created study session."""
    id: str  # UUID as string
    title: str
    studyContent: str
    fileContent: Optional[str] = None  # Original file (base64)
    fileType: Optional[str] = None  # File type: pdf, pptx, docx, txt
    pdfContent: Optional[str] = None  # Converted PDF for PPTX files (base64)
    extractedTopics: List[TopicSchema]
    progress: int
    topics: int
    hasFullStudy: bool
    hasSpeedRun: bool
    createdAt: Optional[int] = None  # Unix timestamp in milliseconds


@router.post("/analyze-content", response_model=ContentAnalysisResponse)
@limiter.limit("10/minute")
async def analyze_content(
    request: Request,
    data: AnalyzeContentRequest,
    current_user: User = Depends(get_current_active_user),
):
    """
    Analyze content and provide recommendations for topics and questions.

    This endpoint analyzes the uploaded content and returns:
    - Word count and estimated reading time
    - Recommended number of topics
    - Recommended questions per topic
    - Complexity score

    Rate Limits:
        - 10 requests per minute per user
    """
    try:
        # Extract text from content (handles Word docs, PDFs, plain text)
        extracted_text = extract_text_from_content(data.content)

        if not extracted_text or len(extracted_text.strip()) < 50:
            raise HTTPException(
                status_code=400,
                detail="Content is too short or empty. Please provide substantial study material (at least 50 characters)."
            )

        # Analyze content complexity
        analysis = analyze_content_complexity(extracted_text)

        # Generate a brief summary (first 200 characters)
        content_preview = extracted_text[:200].strip()
        if len(extracted_text) > 200:
            content_preview += "..."

        return ContentAnalysisResponse(
            word_count=analysis['word_count'],
            estimated_reading_time=analysis['estimated_reading_time'],
            recommended_topics=analysis['recommended_topics'],
            recommended_questions=analysis['recommended_questions'],
            complexity_score=analysis['complexity_score'],
            content_summary=content_preview
        )

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to analyze content: {str(e)}"
        )


@router.post("/create-with-ai", response_model=CreateStudySessionResponse)
@limiter.limit("5/minute")
async def create_study_session_with_ai(
    request: Request,
    data: CreateStudySessionRequest,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db),
):
    """
    Create a new study session and generate topics/questions using AI.

    This endpoint:
    1. Analyzes the provided content
    2. Extracts major topics using AI (dynamically based on content)
    3. Generates quiz questions for each topic
    4. Saves everything to the database

    Rate Limits:
        - 5 requests per minute per user
    """
    try:
        # Validate file size before processing
        # Note: Base64 encoding increases file size by ~33%, so 35MB raw = ~47MB encoded
        MAX_FILE_SIZE_MB = 35
        MAX_FILE_SIZE_BYTES = MAX_FILE_SIZE_MB * 1024 * 1024
        content_size = len(data.content)

        if content_size > MAX_FILE_SIZE_BYTES:
            logger.error(f"❌ File too large: {content_size / (1024*1024):.1f}MB (max: {MAX_FILE_SIZE_MB}MB)")
            raise HTTPException(
                status_code=413,
                detail=f"File size ({content_size / (1024*1024):.1f}MB) exceeds maximum allowed size ({MAX_FILE_SIZE_MB}MB). Please use a smaller file or split it into multiple documents."
            )

        logger.info(f"📁 Processing file: {content_size / (1024*1024):.1f}MB")

        # Extract text and detect file type
        extracted_text, file_type, file_content = detect_file_type_and_extract(data.content)

        if not extracted_text or len(extracted_text.strip()) < 50:
            raise HTTPException(
                status_code=400,
                detail="Content is too short or empty. Please provide substantial study material (at least 50 characters)."
            )

        # Convert PPTX to PDF for better rendering in Speed Run mode
        pdf_content = None
        if file_type == 'pptx':
            logger.info("🔄 Converting PowerPoint to PDF for enhanced rendering...")
            pdf_content = convert_pptx_to_pdf(file_content)
            if pdf_content:
                logger.info("✅ PowerPoint successfully converted to PDF")
            else:
                logger.warning("⚠️  PowerPoint to PDF conversion failed - will use extracted text fallback")

        # Check document size and use chunking for large documents
        estimated_tokens = len(extracted_text) // 4  # Rough estimate: 1 token ≈ 4 characters
        logger.info(f"📊 Document size: {len(extracted_text):,} chars, ~{estimated_tokens:,} estimated tokens")

        # Claude 3.5 Haiku has 200k token context window
        # Use chunking for documents that would exceed safe limits
        # Reduced significantly to account for prompt overhead (instructions, subtopics list, output, etc.)
        CHUNK_SIZE_TOKENS = 20000  # 20k tokens per chunk (~80k chars) - conservative to prevent "too large" errors
        OVERLAP_TOKENS = 2000  # 2k token overlap between chunks for context preservation

        chunk_size_chars = CHUNK_SIZE_TOKENS * 4
        overlap_chars = OVERLAP_TOKENS * 4

        # Split document into chunks if necessary
        document_chunks = []
        if estimated_tokens > CHUNK_SIZE_TOKENS:
            logger.info(f"📚 Large document detected. Splitting into chunks of ~{CHUNK_SIZE_TOKENS:,} tokens with {OVERLAP_TOKENS:,} token overlap...")

            # Split into overlapping chunks
            current_pos = 0
            chunk_num = 1
            while current_pos < len(extracted_text):
                end_pos = min(current_pos + chunk_size_chars, len(extracted_text))
                chunk = extracted_text[current_pos:end_pos]
                document_chunks.append(chunk)
                logger.info(f"  📄 Chunk {chunk_num}: {len(chunk):,} chars (~{len(chunk)//4:,} tokens)")

                # Move forward with overlap
                current_pos = end_pos - overlap_chars
                if current_pos >= len(extracted_text) - overlap_chars:
                    break  # Last chunk
                chunk_num += 1

            logger.info(f"✅ Created {len(document_chunks)} chunks for processing")
        else:
            # Document fits in one chunk
            document_chunks = [extracted_text]
            logger.info(f"✅ Document fits in single chunk, no splitting needed")

        # Analyze content to get smart recommendations (use first chunk for analysis)
        analysis = analyze_content_complexity(document_chunks[0])

        # Progressive loading: For large documents (>5000 words), start with fewer topics
        is_large_doc = analysis['word_count'] > 5000
        initial_topics = data.num_topics

        if data.progressive_load and is_large_doc:
            # Start with only 2-3 categories and 4-6 subtopics for quick initial load
            initial_topics = min(6, data.num_topics)
            logger.info(f"📚 Large document detected ({analysis['word_count']} words). Using progressive load: {initial_topics} initial topics")

        # Calculate number of categories based on topics to generate
        # For 2-5 topics: 2 categories
        # For 6-10 topics: 3 categories
        # For 11-15 topics: 4 categories
        # For 16-20 topics: 5 categories
        num_categories = max(2, min(5, (initial_topics + 3) // 4))
        subtopics_per_category = max(1, initial_topics // num_categories)  # Allow single subtopic per category

        # Initialize AI client (prefer Claude Haiku for speed, fallback to DeepSeek)
        use_claude = bool(settings.ANTHROPIC_API_KEY)
        anthropic_client = None
        deepseek_client = None

        if use_claude:
            anthropic_client = Anthropic(api_key=settings.ANTHROPIC_API_KEY)
            logger.info("🚀 Using Claude Haiku for fast generation")

        # Always initialize DeepSeek as fallback
        if settings.DEEPSEEK_API_KEY:
            deepseek_client = OpenAI(
                api_key=settings.DEEPSEEK_API_KEY,
                base_url="https://api.deepseek.com"
            )
            if not use_claude:
                logger.info("⏱️ Using DeepSeek (consider adding ANTHROPIC_API_KEY for 10x speed)")

        # Step 1: Extract topics from content with hierarchical structure (DYNAMIC PROMPT)
        topics_prompt = f"""Analyze this study material and organize it into a clear, focused hierarchical structure.

Study Material:
{extracted_text}

Content Analysis:
- Word Count: {analysis['word_count']}
- Complexity Score: {analysis['complexity_score']}
- Estimated Reading Time: {analysis['estimated_reading_time']} minutes

Requirements:
1. Create approximately {num_categories} major categories that organize the content at a high level
2. Within each category, break down into focused subtopics (2-3 levels maximum)
3. Prioritize QUALITY over excessive nesting - only create subtopics when there's sufficient content to generate meaningful questions
4. Each leaf node should be a focused concept that can support 5-15 quality questions
5. Provide clear titles and brief descriptions at ALL levels
6. Organize logically (foundational concepts first, building to advanced topics)
7. Create approximately {initial_topics} total LEAF topics across all categories that will actually have questions
8. Avoid creating unnecessary intermediate categories - go directly to testable concepts when possible
9. Focus on core concepts and foundational topics first
10. Remember: The goal is quality questions, not deep nesting

IMPORTANT: Keep the structure simple and focused. Only nest 2-3 levels deep. Empty subtopics with no questions provide no value.

Return ONLY a valid JSON object in this EXACT format (keep to 2-3 levels maximum):
{{
  "categories": [
    {{
      "title": "Category Title (Level 1)",
      "description": "Brief description of this category",
      "subtopics": [
        {{
          "title": "Subtopic Title (Level 2)",
          "description": "Brief description",
          "subtopics": [
            {{
              "title": "Focused Topic (Level 3 - Leaf)",
              "description": "Brief description of testable concept",
              "subtopics": []
            }}
          ]
        }},
        {{
          "title": "Another Subtopic (Level 2 - Leaf)",
          "description": "Can also be a leaf node if content is focused enough",
          "subtopics": []
        }}
      ]
    }}
  ]
}}

Note: An EMPTY subtopics array [] means this is a LEAF NODE that will have questions generated for it. Keep nesting to 2-3 levels maximum."""

        # Call AI to extract topics (with automatic fallback to DeepSeek if Claude fails)
        topics_text = None
        if use_claude and anthropic_client:
            try:
                topics_response = anthropic_client.messages.create(
                    model="claude-3-5-haiku-20241022",
                    max_tokens=2048,
                    temperature=0.7,
                    messages=[{"role": "user", "content": topics_prompt}]
                )
                topics_text = topics_response.content[0].text
            except Exception as claude_error:
                logger.warning(f"⚠️ Claude API failed: {str(claude_error)}")
                if deepseek_client:
                    logger.info("🔄 Falling back to DeepSeek...")
                    use_claude = False  # Switch to DeepSeek for remaining calls
                else:
                    raise HTTPException(
                        status_code=500,
                        detail=f"Claude API failed and no DeepSeek fallback available: {str(claude_error)}"
                    )

        if not topics_text and deepseek_client:
            topics_response = deepseek_client.chat.completions.create(
                model="deepseek-chat",
                max_tokens=2048,
                temperature=0.7,
                messages=[{"role": "user", "content": topics_prompt}]
            )
            topics_text = topics_response.choices[0].message.content

        if not topics_text:
            raise HTTPException(
                status_code=500,
                detail="No AI provider available. Please configure ANTHROPIC_API_KEY or DEEPSEEK_API_KEY."
            )

        # Parse hierarchical topics
        try:
            start_idx = topics_text.find('{')
            end_idx = topics_text.rfind('}') + 1
            topics_json = json.loads(topics_text[start_idx:end_idx])
            categories_data = topics_json["categories"]
        except (json.JSONDecodeError, KeyError, ValueError) as e:
            raise HTTPException(status_code=500, detail=f"Failed to parse topics: {str(e)}")

        # Count total subtopics for the session
        total_subtopics = sum(len(cat.get("subtopics", [])) for cat in categories_data)

        # Generate a smart title if the provided title is generic (contains "Study Session" and a date)
        session_title = data.title
        if "Study Session" in session_title and any(char.isdigit() for char in session_title):
            # Use the first category as the title for better organization
            if categories_data and len(categories_data) > 0:
                first_category = categories_data[0]["title"]
                # If multiple categories, add a subtitle
                if len(categories_data) > 1:
                    session_title = f"{first_category} + {len(categories_data) - 1} more"
                else:
                    session_title = first_category

        # Step 2: Create study session
        study_session = StudySession(
            user_id=current_user.id,
            title=session_title,
            topic=categories_data[0]["title"] if categories_data else "General Study",
            study_content=extracted_text,  # Store the extracted text
            file_content=file_content,  # Store original file (base64)
            file_type=file_type,  # Store file type (pdf, pptx, docx, txt)
            pdf_content=pdf_content,  # Store converted PDF for PPTX files (base64)
            topics_count=total_subtopics,
            has_full_study=True,
            has_speed_run=True,
            status="in_progress"
        )
        db.add(study_session)
        db.flush()  # Get the session ID

        # Step 3: Create categories and subtopics first, then batch generate ALL questions
        all_topics = []
        subtopic_map = {}  # Map to track subtopics for batch question assignment
        overall_idx = 0

        # Recursive helper function to create topics at any depth
        def create_topics_recursive(
            topic_data: dict,
            parent_topic_id: Optional[int],
            parent_schema: Optional[TopicSchema],
            path: str,  # e.g., "0-1-2" for tracking hierarchy
            order_index: int,
            current_depth: int = 0
        ) -> Optional[TopicSchema]:
            """
            Recursively create topics up to MAX_DEPTH levels.
            Returns the schema for this topic (or None if it's a pure container).
            """
            MAX_DEPTH = 3  # Limit to 3 levels: Category -> Subtopic -> Leaf (or Category -> Leaf)

            subtopics_data = topic_data.get("subtopics", [])

            # Force leaf node if we've reached max depth, regardless of subtopics in data
            if current_depth >= MAX_DEPTH:
                subtopics_data = []
                if topic_data.get("subtopics"):
                    logger.warning(f"⚠️ Max depth {MAX_DEPTH} reached for topic '{topic_data['title']}' - forcing as leaf node")

            has_children = len(subtopics_data) > 0
            is_leaf = not has_children

            # Create topic in database
            topic = Topic(
                study_session_id=study_session.id,
                parent_topic_id=parent_topic_id,
                title=topic_data["title"],
                description=topic_data.get("description", ""),
                order_index=order_index,
                is_category=has_children  # Non-leaf nodes are categories
            )
            db.add(topic)
            db.flush()

            # Create schema
            topic_schema = TopicSchema(
                id=f"topic-{path}",
                db_id=topic.id,
                title=topic_data["title"],
                description=topic_data.get("description", ""),
                questions=[],
                completed=False,
                score=None,
                currentQuestionIndex=0,
                isCategory=has_children,
                parentTopicId=parent_schema.id if parent_schema else None,
                subtopics=[]
            )

            # Add ALL topics (both categories and leaf nodes) to subtopic_map for question generation
            # Categories will get overview/synthesis questions, leaf nodes get specific questions
            subtopic_map[path] = {
                "topic": topic,
                "topic_data": topic_data,
                "schema": topic_schema,
                "parent_schema": parent_schema,
                "is_category": has_children,
                "is_leaf": is_leaf
            }

            # Recursively create children
            if has_children:
                for child_idx, child_data in enumerate(subtopics_data):
                    child_path = f"{path}-{child_idx}"
                    child_schema = create_topics_recursive(
                        child_data,
                        topic.id,
                        topic_schema,
                        child_path,
                        child_idx,
                        current_depth + 1  # Increment depth for children
                    )
                    if child_schema:
                        topic_schema.subtopics.append(child_schema)

            return topic_schema

        # First pass: Create all categories and their nested subtopics recursively
        for cat_idx, category_data in enumerate(categories_data):
            category_schema = create_topics_recursive(
                category_data,
                parent_topic_id=None,
                parent_schema=None,
                path=str(cat_idx),
                order_index=cat_idx
            )
            if category_schema:
                all_topics.append(category_schema)

        # Step 4: Generate questions by processing each document chunk
        # For large documents, this processes multiple chunks separately and merges results
        logger.info(f"📡 Processing {len(document_chunks)} document chunk(s) to generate questions...")

        # Generate questions for ALL subtopics during initial creation
        # This ensures maximum question coverage from the uploaded document

        # Get list of all subtopic keys
        all_subtopic_keys = list(subtopic_map.keys())

        logger.info(f"📚 Generating questions for ALL {len(all_subtopic_keys)} subtopics to maximize coverage")
        logger.info(f"  Subtopics: {all_subtopic_keys}")

        # BATCH SUBTOPICS: Process in groups to avoid AI refusing due to response size
        # With 2 subtopics per batch, max ~60 questions per request (2 * 30) - ensures ALL topics get questions
        SUBTOPICS_PER_BATCH = 2  # Reduced to ensure AI generates questions for ALL topics
        subtopic_batches = []

        # Split subtopics into batches
        for i in range(0, len(all_subtopic_keys), SUBTOPICS_PER_BATCH):
            batch_keys = all_subtopic_keys[i:i + SUBTOPICS_PER_BATCH]
            subtopic_batches.append(batch_keys)

        logger.info(f"📦 Split {len(all_subtopic_keys)} subtopics into {len(subtopic_batches)} batches of up to {SUBTOPICS_PER_BATCH}")

        # Collect questions from all chunks
        all_chunk_questions = {}  # {subtopic_key: [questions]}

        # Process each chunk with batched subtopics
        for chunk_idx, chunk_text in enumerate(document_chunks, 1):
            logger.info(f"📄 Processing chunk {chunk_idx}/{len(document_chunks)}...")

            # Process each batch of subtopics for this chunk
            for batch_num, batch_keys in enumerate(subtopic_batches, 1):
                logger.info(f"  📦 Batch {batch_num}/{len(subtopic_batches)}: Processing {len(batch_keys)} subtopics")

                # Build subtopics list for this batch only
                subtopics_list = ""
                for subtopic_key in batch_keys:
                    subtopic_info = subtopic_map[subtopic_key]
                    topic_data = subtopic_info["topic_data"]
                    is_category = subtopic_info.get("is_category", False)
                    is_leaf = subtopic_info.get("is_leaf", True)

                    # Simple format that works for any nesting depth
                    subtopics_list += f"\n[Topic {subtopic_key}]\n"
                    subtopics_list += f"Title: {topic_data['title']}\n"
                    subtopics_list += f"Description: {topic_data.get('description', '')}\n"
                    subtopics_list += f"Type: {'CATEGORY (overview/synthesis questions)' if is_category else 'SPECIFIC TOPIC (detailed questions)'}\n"

                # Build prompt for this chunk and subtopic batch
                batch_prompt = f"""Generate TRICKY and CHALLENGING multiple-choice questions for EACH of the following topics from the study material.

CRITICAL: Generate high-quality, comprehensive questions:
- Extract the MOST IMPORTANT testable concepts, facts, principles, definitions, and examples from the material
- Aim for 15-20 high-quality questions per topic (quality over quantity)
- Break down key concepts into multiple questions from different angles
- Test each concept in multiple ways: definition, application, comparison, analysis, synthesis, evaluation
- Create questions for the most important content in the study material
- Generate questions for ALL listed topics below - if the material is light on a topic, use the topic title and description to create questions
- Even if a topic is only mentioned briefly, create comprehensive questions based on what IS mentioned
- For topics not directly covered in the chunk, create questions based on the topic title/description and related content
- Focus on core concepts and ensure complete coverage across all topics

QUESTION TYPES BASED ON TOPIC TYPE:
1. For CATEGORY topics (overview/synthesis questions):
   - Test overall understanding of the category and how subtopics relate
   - Ask synthesis questions that integrate concepts from multiple subtopics
   - Include comparison questions between different subtopics within the category
   - Test the big picture and overarching principles

2. For SPECIFIC TOPIC (detailed questions):
   - Test deep, specific knowledge about that particular topic
   - Include detailed questions about definitions, principles, and mechanisms
   - Ask application questions specific to that topic

3. For ALL topics (CRITICAL - include these):
   - EXAMPLE-BASED QUESTIONS: "Give an example of...", "Which scenario demonstrates...", "Which of the following is an example of..."
   - Application questions that require applying concepts to real scenarios
   - Questions that test practical understanding through examples

DIFFICULTY LEVEL: CHALLENGING
- Make questions that require DEEP analysis and critical thinking
- Use tricky distractors that would fool someone who only skimmed the material
- Test subtle distinctions and nuanced understanding
- Require application of concepts, not just memorization
- Include "all of the above" or "none of the above" when appropriate
- Use comparative questions (e.g., "Which is the PRIMARY..." "What is the MAIN difference...")
- Create questions that test WHY and HOW, not just WHAT
- For every concept, include at least one EXAMPLE-BASED question

Study Material (Chunk {chunk_idx} of {len(document_chunks)}):
{chunk_text}

TOPICS TO COVER ({len(batch_keys)} topics in this batch):
{subtopics_list}

Requirements:
1. Generate 15-20 high-quality questions for EACH topic (total ~30-40 questions for this batch)
2. Extract the MOST IMPORTANT testable information from the study material
3. NO DUPLICATES - each question must test a unique concept or angle
4. For EACH topic, include multiple EXAMPLE-BASED questions (e.g., "Which scenario is an example of operant conditioning?")
5. Each question must have exactly 4 PLAUSIBLE options (all should seem correct to someone who doesn't understand deeply)
6. Questions should be TRICKY and CHALLENGING - test deep understanding and critical thinking
7. Distractors should be subtle and based on common misconceptions
8. Provide detailed explanations that explain why the correct answer is right AND why the distractors are wrong
9. For EACH question, include the EXACT source text from the study material with FULL CONTEXT
10. Source text should include the complete sentence(s) or paragraph that contains the answer
11. Include enough surrounding context (2-4 sentences) so students can easily locate it in their document
12. The sourceText must be VERBATIM from the study material - copy it EXACTLY as it appears
13. For PDF documents, estimate the page number where this content appears (if this is chunk {chunk_idx} of {len(document_chunks)}, estimate accordingly)
14. Return ONLY valid JSON - NO MARKDOWN, NO CODE BLOCKS, NO EXTRA TEXT

JSON FORMATTING RULES (CRITICAL):
- Use double quotes (") for all strings, not single quotes
- Escape special characters: \" for quotes, \\ for backslashes, \n for newlines
- Do NOT use trailing commas after the last item in arrays or objects
- Ensure all brackets and braces are properly closed
- Numbers for correctAnswer should be integers (0, 1, 2, 3), not strings
- Do NOT include comments in the JSON

IMPORTANT INSTRUCTIONS:
- Do NOT ask clarifying questions
- Do NOT say "I understand" or provide explanations
- Do NOT add any text before or after the JSON
- START your response IMMEDIATELY with the opening brace: {{
- Generate the complete JSON response now

GOAL: Create 15-20 comprehensive questions per topic. Focus on QUALITY and coverage of core concepts. Include example-based questions for key concepts.

CRITICAL REQUIREMENTS:
- You MUST generate questions for EVERY SINGLE topic key listed above - NO EXCEPTIONS
- If a topic is listed, it MUST appear in your JSON response with questions
- Missing even ONE topic key will result in incomplete learning coverage
- Each topic MUST have 15-20 questions (aim for 20 when possible)

Return in this EXACT format (use topic keys EXACTLY as shown above - EVERY topic listed must be in the response):
{{
  "subtopics": {{
    "0": {{
      "questions": [
        {{
          "question": "Question text?",
          "options": ["Option A", "Option B", "Option C", "Option D"],
          "correctAnswer": 0,
          "explanation": "Why this answer is correct",
          "sourceText": "The complete sentence or paragraph from the study material with context.",
          "sourcePage": null
        }},
        ... (10-30+ questions for topic "0")
      ]
    }},
    "0-1": {{
      "questions": [
        {{question object}},
        ... (10-30+ questions for topic "0-1")
      ]
    }},
    ... (MUST include ALL topic keys from the list above)
  }}
}}

REMINDER: The response MUST include questions for ALL {len(batch_keys)} topics listed above. Double-check that every topic key appears in your JSON response."""

                # Make API call for this batch
                prompt_tokens = len(batch_prompt) // 4  # Rough estimate
                logger.info(f"    📊 Batch {batch_num} prompt length: {len(batch_prompt):,} characters (~{prompt_tokens:,} tokens)")

                # Note: Removed hard prompt size check - let AI handle it gracefully with error recovery below
                # With reduced chunk size (20k) and batch size (3 subtopics), prompts should be safe

                batch_text = None
                try:
                    if use_claude and anthropic_client:
                        try:
                            batch_response = anthropic_client.messages.create(
                                model="claude-3-5-haiku-20241022",
                                max_tokens=8192,  # Maximum output tokens for Claude 3.5 Haiku
                                temperature=0.7,
                                messages=[{"role": "user", "content": batch_prompt}]
                            )
                            batch_text = batch_response.content[0].text
                        except Exception as claude_error:
                            logger.warning(f"⚠️ Claude API failed for chunk {chunk_idx} batch {batch_num}: {str(claude_error)}")
                            if deepseek_client:
                                logger.info(f"🔄 Falling back to DeepSeek for chunk {chunk_idx} batch {batch_num}...")
                                use_claude = False  # Switch to DeepSeek for remaining calls
                            else:
                                raise  # Re-raise if no fallback available

                    if not batch_text and deepseek_client:
                        batch_response = deepseek_client.chat.completions.create(
                            model="deepseek-chat",
                            max_tokens=64000,  # DeepSeek supports larger output
                            temperature=0.7,
                            messages=[{"role": "user", "content": batch_prompt}]
                        )
                        batch_text = batch_response.choices[0].message.content

                    if not batch_text:
                        raise Exception("No AI provider available")

                except Exception as api_error:
                    logger.warning(f"⚠️ AI API call failed for chunk {chunk_idx} batch {batch_num}: {type(api_error).__name__}: {str(api_error)}")
                    # Check if it's a context length error
                    error_msg = str(api_error).lower()
                    if any(keyword in error_msg for keyword in ['context', 'token', 'too long', 'maximum', 'limit']):
                        logger.warning(f"⚠️ Skipping chunk {chunk_idx} batch {batch_num} - too large. Continuing with other batches...")
                        continue  # Skip this batch and continue with next one
                    else:
                        # For non-context errors, still fail (e.g., auth issues, network errors)
                        raise HTTPException(
                            status_code=500,
                            detail=f"AI API error on chunk {chunk_idx} batch {batch_num}: {str(api_error)}"
                        )

                # Log AI response details
                logger.info(f"    📨 Batch {batch_num} - Received AI response, length: {len(batch_text)} characters")

                # Parse batch response with detailed error logging
                try:
                    start_idx = batch_text.find('{')
                    end_idx = batch_text.rfind('}') + 1

                    if start_idx == -1 or end_idx == 0:
                        logger.error(f"❌ Chunk {chunk_idx} Batch {batch_num} - No JSON found in AI response")
                        logger.error(f"❌ AI returned: {batch_text[:500]}...")  # Log first 500 chars
                        chunk_questions = {}
                    else:
                        json_str = batch_text[start_idx:end_idx]
                        batch_json = json.loads(json_str)
                        chunk_questions = batch_json.get("subtopics", {})

                        logger.info(f"    ✅ Batch {batch_num} - Parsed {len(chunk_questions)} subtopics")
                        for key in chunk_questions.keys():
                            q_count = len(chunk_questions[key].get("questions", []))
                            if q_count > 0:
                                logger.info(f"      - Subtopic {key}: {q_count} questions")

                        # Validate that ALL topics in this batch got questions
                        missing_topics = [key for key in batch_keys if key not in chunk_questions or not chunk_questions[key].get("questions")]
                        if missing_topics:
                            logger.warning(f"    ⚠️ Batch {batch_num} - AI did NOT generate questions for {len(missing_topics)} topics: {missing_topics}")
                            logger.warning(f"    ⚠️ This may result in incomplete coverage. Topics requested: {batch_keys}")
                        else:
                            logger.info(f"    ✅ Batch {batch_num} - ALL {len(batch_keys)} topics have questions")

                except (json.JSONDecodeError, KeyError, ValueError) as e:
                    logger.error(f"❌ Chunk {chunk_idx} Batch {batch_num} - Failed to parse questions: {e}")
                    chunk_questions = {}

                # Merge questions from this batch into all_chunk_questions
                for subtopic_key, subtopic_data in chunk_questions.items():
                    questions = subtopic_data.get("questions", [])
                    if questions:
                        if subtopic_key not in all_chunk_questions:
                            all_chunk_questions[subtopic_key] = []
                        all_chunk_questions[subtopic_key].extend(questions)
                        logger.debug(f"    🔄 Added {len(questions)} questions for subtopic {subtopic_key}")

        # Log merged results
        logger.info(f"🎯 Merging complete - Total subtopics with questions: {len(all_chunk_questions)}")
        for key, questions in all_chunk_questions.items():
            logger.info(f"  - Subtopic {key}: {len(questions)} total questions from all chunks")

        # Use merged questions as the final subtopics_questions
        subtopics_questions = {key: {"questions": questions} for key, questions in all_chunk_questions.items()}

        # Step 5: Assign questions to subtopics
        question_counter = 0
        logger.info(f"🔄 Assigning questions to {len(subtopic_map)} subtopics...")
        logger.debug(f"🗂️ Subtopic keys in map: {list(subtopic_map.keys())}")
        logger.debug(f"🗂️ Subtopic keys in AI response: {list(subtopics_questions.keys())}")

        for subtopic_key, subtopic_info in subtopic_map.items():
            subtopic = subtopic_info["topic"]
            subtopic_schema = subtopic_info["schema"]
            # Get parent schema (either category_schema or subtopic_schema)
            parent_schema = subtopic_info.get("category_schema") or subtopic_info.get("subtopic_schema")

            # Get questions for this subtopic from batch response
            questions_data = subtopics_questions.get(subtopic_key, {}).get("questions", [])

            # Skip if no questions generated
            if not questions_data:
                # Expected to have questions for all subtopics
                logger.warning(f"⚠️ No questions generated for subtopic '{subtopic_key}' ('{subtopic.title}') - SKIPPING (may lack relevant content in document)")
                # If this is a sub-subtopic without questions, still add it to parent
                if "subtopic_schema" in subtopic_info:
                    parent_schema.subtopics.append(subtopic_schema)
                continue
            else:
                logger.info(f"✅ Found {len(questions_data)} questions for subtopic '{subtopic_key}' ('{subtopic.title}')")

            # Save ALL questions to database (no limit - variable per topic)
            # Get existing questions for this topic to check for duplicates
            existing_questions = db.query(Question).filter(Question.topic_id == subtopic.id).all()
            existing_question_texts = {q.question.lower().strip() for q in existing_questions}

            questions_list = []
            duplicates_skipped = 0

            for q_idx, q_data in enumerate(questions_data):
                question_text = q_data.get("question", f"Question {q_idx+1}")

                # Check for duplicate questions (case-insensitive)
                if question_text.lower().strip() in existing_question_texts:
                    logger.debug(f"⏭️ Skipping duplicate question: {question_text[:50]}...")
                    duplicates_skipped += 1
                    continue

                # Ensure options is a list
                options = q_data.get("options", ["A", "B", "C", "D"])
                if isinstance(options, str):
                    try:
                        options = json.loads(options)
                    except json.JSONDecodeError:
                        options = ["Option A", "Option B", "Option C", "Option D"]

                question = Question(
                    topic_id=subtopic.id,
                    question=question_text,
                    options=options,
                    correct_answer=q_data.get("correctAnswer", 0),
                    explanation=q_data.get("explanation", ""),
                    source_text=q_data.get("sourceText"),
                    source_page=q_data.get("sourcePage"),
                    order_index=len(questions_list)  # Use actual index in list
                )
                db.add(question)

                # Add to existing set to catch duplicates within this batch
                existing_question_texts.add(question_text.lower().strip())

                questions_list.append(QuestionSchema(
                    id=f"q-{question_counter}",
                    question=q_data.get("question", f"Question {q_idx+1}"),
                    options=options,
                    correctAnswer=q_data.get("correctAnswer", 0),
                    explanation=q_data.get("explanation", ""),
                    sourceText=q_data.get("sourceText"),
                    sourcePage=q_data.get("sourcePage")
                ))
                question_counter += 1

            # Log duplicate detection results
            if duplicates_skipped > 0:
                logger.info(f"  ⏭️ Skipped {duplicates_skipped} duplicate questions for '{subtopic.title}'")

            # Update subtopic schema with questions
            subtopic_schema.questions = questions_list
            # Add to parent (either category or parent subtopic)
            if parent_schema:
                parent_schema.subtopics.append(subtopic_schema)

        # Validate that questions were generated for a reasonable number of subtopics
        subtopics_with_questions = len([k for k, v in subtopics_questions.items() if v.get("questions")])
        total_subtopics = len(subtopic_map)

        if question_counter == 0:
            logger.error("❌ FATAL: No questions were generated for any subtopic! AI generation completely failed.")
            db.rollback()
            raise HTTPException(
                status_code=500,
                detail="Failed to generate questions. The AI did not return any valid questions. This may be due to document format or content issues. Please try:\n1. A different document\n2. Splitting the document into smaller files\n3. Converting to PDF format if using Word/PowerPoint"
            )

        # Enforce minimum question count for good study sessions
        MIN_QUESTIONS_REQUIRED = 10  # Lowered for nested subtopics that distribute questions across more leaf nodes
        if question_counter < MIN_QUESTIONS_REQUIRED:
            logger.error(f"❌ Only {question_counter} questions generated (minimum: {MIN_QUESTIONS_REQUIRED})")
            db.rollback()
            raise HTTPException(
                status_code=500,
                detail=f"Not enough questions generated ({question_counter}/{MIN_QUESTIONS_REQUIRED} required). The document may be too short, too complex, or in an unsupported format. Please try:\n1. A longer or more detailed document\n2. Converting to PDF format\n3. Checking that the content is educational material"
            )

        # Log coverage statistics
        coverage_percent = (subtopics_with_questions / total_subtopics * 100) if total_subtopics > 0 else 0
        logger.info(f"📊 Question generation coverage: {subtopics_with_questions}/{total_subtopics} subtopics ({coverage_percent:.1f}%)")
        logger.info(f"✅ Successfully generated {question_counter} total questions across {subtopics_with_questions} subtopics")

        # Commit all changes
        db.commit()
        db.refresh(study_session)

        return CreateStudySessionResponse(
            id=str(study_session.id),  # Convert UUID to string
            title=study_session.title,
            studyContent=study_session.study_content,
            fileContent=study_session.file_content,
            fileType=study_session.file_type,
            pdfContent=study_session.pdf_content,  # Converted PDF for PPTX files
            extractedTopics=all_topics,
            progress=0,
            topics=len(all_topics),
            hasFullStudy=True,
            hasSpeedRun=True,
            createdAt=int(study_session.created_at.timestamp() * 1000) if study_session.created_at else None
        )

    except HTTPException:
        # Re-raise HTTP exceptions (like 413 for file size)
        db.rollback()
        raise
    except Exception as api_error:
        logger.error(f"❌ Error in create_study_session_with_ai: {type(api_error).__name__}: {str(api_error)}")
        logger.error(f"❌ Full error: {repr(api_error)}", exc_info=True)
        if "openai" in str(type(api_error).__module__):
            db.rollback()
            raise HTTPException(status_code=500, detail=f"DeepSeek API error: {str(api_error)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Failed to create study session: {str(api_error)}")


@router.get("/{session_id}", response_model=CreateStudySessionResponse)
async def get_study_session(
    session_id: str,  # UUID as string
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db),
):
    """
    Get a study session with all its topics and questions.

    Returns the complete session data including the hierarchical topic structure
    with all questions, allowing users to resume their study progress.
    """
    # Validate UUID format
    try:
        uuid_obj = uuid.UUID(session_id)
    except ValueError:
        logger.error(f"❌ Invalid session ID format: {session_id} (expected UUID)")
        raise HTTPException(
            status_code=400,
            detail=f"Invalid session ID format. Expected UUID, got: '{session_id}'. This may be an old session that is no longer compatible."
        )

    # Fetch session with eager loading of topics and questions
    session = db.query(StudySession).filter(
        StudySession.id == uuid_obj,
        StudySession.user_id == current_user.id
    ).first()

    if not session:
        raise HTTPException(status_code=404, detail="Study session not found")

    # Fetch all topics for this session
    all_topics = db.query(Topic).filter(
        Topic.study_session_id == uuid_obj
    ).order_by(Topic.order_index).all()

    # Build hierarchical structure
    categories = [t for t in all_topics if t.is_category and t.parent_topic_id is None]

    result_topics = []

    for cat_idx, category in enumerate(categories):
        # Get subtopics for this category
        subtopics = [t for t in all_topics if t.parent_topic_id == category.id]

        category_schema = TopicSchema(
            id=f"category-{cat_idx+1}",
            db_id=category.id,  # Include database ID
            title=category.title,
            description=category.description or "",
            isCategory=True,
            parentTopicId=None,
            questions=[],
            subtopics=[]
        )

        for sub_idx, subtopic in enumerate(subtopics):
            # Fetch questions for this subtopic
            questions = db.query(Question).filter(
                Question.topic_id == subtopic.id
            ).order_by(Question.order_index).all()

            questions_list = [
                QuestionSchema(
                    id=f"topic-{sub_idx+1}-q{q.order_index+1}",
                    question=q.question,
                    options=q.options,
                    correctAnswer=q.correct_answer,
                    explanation=q.explanation,
                    sourceText=q.source_text,
                    sourcePage=q.source_page
                )
                for q in questions
            ]

            subtopic_schema = TopicSchema(
                id=f"subtopic-{cat_idx+1}-{sub_idx+1}",
                db_id=subtopic.id,  # Include database ID for progress sync
                title=subtopic.title,
                description=subtopic.description or "",
                questions=questions_list,
                completed=subtopic.completed or False,
                score=subtopic.score,
                currentQuestionIndex=subtopic.current_question_index or 0,
                isCategory=False,
                parentTopicId=f"category-{cat_idx+1}",
                subtopics=[]
            )
            category_schema.subtopics.append(subtopic_schema)

        result_topics.append(category_schema)

    # Calculate progress
    total_subtopics = sum(len(cat.subtopics) for cat in result_topics)
    completed_subtopics = sum(
        sum(1 for st in cat.subtopics if st.completed)
        for cat in result_topics
    )
    progress = int((completed_subtopics / total_subtopics * 100) if total_subtopics > 0 else 0)

    return CreateStudySessionResponse(
        id=str(session.id),  # Convert UUID to string
        title=session.title,
        studyContent=session.study_content or "",
        fileContent=session.file_content,
        fileType=session.file_type,
        pdfContent=session.pdf_content,  # Converted PDF for PPTX files
        extractedTopics=result_topics,
        progress=progress,
        topics=total_subtopics,
        hasFullStudy=session.has_full_study or False,
        hasSpeedRun=session.has_speed_run or False,
        createdAt=int(session.created_at.timestamp() * 1000) if session.created_at else None
    )


@router.post("/{session_id}/generate-more-questions")
async def generate_more_questions(
    session_id: str,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db),
):
    """
    Generate questions for remaining subtopics that don't have questions yet.

    This implements progressive loading - initially only first few subtopics have questions,
    calling this endpoint generates questions for the next batch.
    """
    logger.info(f"📚 Generating more questions for session {session_id}")

    # Validate UUID
    try:
        uuid_obj = uuid.UUID(session_id)
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid session ID format")

    # Fetch session
    session = db.query(StudySession).filter(
        StudySession.id == uuid_obj,
        StudySession.user_id == current_user.id
    ).first()

    if not session:
        raise HTTPException(status_code=404, detail="Study session not found")

    # OPTIMIZED: Get all topics with question counts in a single aggregated query
    # Subquery to count questions per topic
    question_counts_subquery = db.query(
        Question.topic_id,
        func.count(Question.id).label('question_count')
    ).group_by(Question.topic_id).subquery()

    # Get topics with their question counts
    topics_with_counts = db.query(
        Topic,
        func.coalesce(question_counts_subquery.c.question_count, 0).label('question_count')
    ).outerjoin(
        question_counts_subquery,
        Topic.id == question_counts_subquery.c.topic_id
    ).filter(
        Topic.study_session_id == uuid_obj
    ).order_by(Topic.order_index).all()

    # Find topics without questions using pre-loaded counts
    subtopics_without_questions = [
        topic for topic, count in topics_with_counts if count == 0
    ]

    if not subtopics_without_questions:
        logger.info(f"✅ All subtopics already have questions for session {session_id}")
        return {"message": "All subtopics already have questions", "generated": 0}

    logger.info(f"📊 Found {len(subtopics_without_questions)} topics without questions")

    # Generate questions for next batch (2 topics at a time for better coverage)
    BATCH_SIZE = 2
    next_batch = subtopics_without_questions[:BATCH_SIZE]

    logger.info(f"🔄 Generating questions for next {len(next_batch)} subtopics...")

    # Extract text from stored file content
    if not session.file_content:
        raise HTTPException(status_code=400, detail="No file content available for this session")

    extracted_text, _, _ = detect_file_type_and_extract(session.file_content)

    # Initialize AI client
    use_claude = bool(settings.ANTHROPIC_API_KEY)
    if use_claude:
        anthropic_client = Anthropic(api_key=settings.ANTHROPIC_API_KEY)
        logger.info("🚀 Using Claude Haiku for question generation")
    else:
        deepseek_client = OpenAI(
            api_key=settings.DEEPSEEK_API_KEY,
            base_url="https://api.deepseek.com"
        )
        logger.info("⏱️ Using DeepSeek")

    # Build prompt for next batch
    subtopics_list = ""
    subtopic_map = {}

    for topic in next_batch:
        # Find parent category
        parent = db.query(Topic).filter(Topic.id == topic.parent_topic_id).first()
        category_title = parent.title if parent else "General"

        # Create mapping key (we'll use topic db_id)
        key = f"topic-{topic.id}"
        subtopic_map[key] = topic

        # Include type information for the AI
        topic_type = "CATEGORY (overview/synthesis questions)" if topic.is_category else "SPECIFIC TOPIC (detailed questions)"

        subtopics_list += f"\n[Topic {key}]\n"
        subtopics_list += f"Title: {topic.title}\n"
        subtopics_list += f"Description: {topic.description or ''}\n"
        subtopics_list += f"Type: {topic_type}\n"
        if parent:
            subtopics_list += f"Parent Category: {category_title}\n"

    # Build prompt
    batch_prompt = f"""Generate TRICKY and CHALLENGING multiple-choice questions for EACH of the following topics from the study material.

CRITICAL: Generate the ABSOLUTE MAXIMUM number of questions possible:
- Extract EVERY testable concept, fact, principle, detail, definition, example, and implication from the material
- DO NOT impose any limits on the number of questions - generate as many as the content supports (aim for 15-30+ questions per topic)
- Break down EVERY concept into multiple questions from different angles
- Test each concept in multiple ways: definition, application, comparison, analysis, synthesis, evaluation
- Create questions for every sentence that contains testable information
- Generate questions for ALL listed topics below - if the material is light on a topic, use the topic title and description to create questions
- Even if a topic is only mentioned briefly, create comprehensive questions based on what IS mentioned
- Continue generating until you have exhausted ALL testable content

QUESTION TYPES BASED ON TOPIC TYPE:
1. For CATEGORY topics (overview/synthesis questions):
   - Test overall understanding of the category and how subtopics relate
   - Ask synthesis questions that integrate concepts from multiple subtopics
   - Include comparison questions between different subtopics within the category
   - Test the big picture and overarching principles

2. For SPECIFIC TOPIC (detailed questions):
   - Test deep, specific knowledge about that particular topic
   - Include detailed questions about definitions, principles, and mechanisms
   - Ask application questions specific to that topic

3. For ALL topics (CRITICAL - include these):
   - EXAMPLE-BASED QUESTIONS: "Give an example of...", "Which scenario demonstrates...", "Which of the following is an example of..."
   - Application questions that require applying concepts to real scenarios
   - Questions that test practical understanding through examples

DIFFICULTY LEVEL: CHALLENGING
- Make questions that require DEEP analysis and critical thinking
- Use tricky distractors that would fool someone who only skimmed the material
- Test subtle distinctions and nuanced understanding
- Require application of concepts, not just memorization
- Include "all of the above" or "none of the above" when appropriate
- Use comparative questions (e.g., "Which is the PRIMARY..." "What is the MAIN difference...")
- Create questions that test WHY and HOW, not just WHAT
- For every concept, include at least one EXAMPLE-BASED question

Study Material:
{extracted_text[:80000]}

TOPICS TO COVER ({len(next_batch)} topics in this batch):
{subtopics_list}

Requirements:
1. Generate 15-20 high-quality questions for EACH topic (total ~30-40 questions for this batch)
2. Extract the MOST IMPORTANT testable information from the study material
3. NO DUPLICATES - each question must test a unique concept or angle
4. For EACH topic, include multiple EXAMPLE-BASED questions (e.g., "Which scenario is an example of...?")
5. Each question must have exactly 4 PLAUSIBLE options (all should seem correct to someone who doesn't understand deeply)
6. Questions should be TRICKY and CHALLENGING - test deep understanding and critical thinking
7. Distractors should be subtle and based on common misconceptions
8. Provide detailed explanations that explain why the correct answer is right AND why the distractors are wrong
9. For EACH question, include the EXACT source text from the study material with FULL CONTEXT
10. Source text should include the complete sentence(s) or paragraph that contains the answer
11. Include enough surrounding context (2-4 sentences) so students can easily locate it in their document
12. The sourceText must be VERBATIM from the study material - copy it EXACTLY as it appears
13. For PDF documents, estimate which section/page the content appears in
14. Return ONLY valid JSON - NO MARKDOWN, NO CODE BLOCKS, NO EXTRA TEXT

JSON FORMATTING RULES (CRITICAL):
- Use double quotes (") for all strings, not single quotes
- Escape special characters: \" for quotes, \\ for backslashes, \n for newlines
- Do NOT use trailing commas after the last item in arrays or objects
- Ensure all brackets and braces are properly closed
- Numbers for correctAnswer should be integers (0, 1, 2, 3), not strings
- Do NOT include comments in the JSON

IMPORTANT INSTRUCTIONS:
- Do NOT ask clarifying questions
- Do NOT say "I understand" or provide explanations
- Do NOT add any text before or after the JSON
- START your response IMMEDIATELY with the opening brace: {{
- Generate the complete JSON response now

GOAL: Create 15-20 comprehensive questions per topic. Focus on QUALITY and coverage of core concepts. Include example-based questions for key concepts.

CRITICAL REQUIREMENTS:
- You MUST generate questions for EVERY SINGLE topic key listed above - NO EXCEPTIONS
- If a topic is listed, it MUST appear in your JSON response with questions
- Missing even ONE topic key will result in incomplete learning coverage
- Each topic MUST have 15-20 questions (aim for 20 when possible)

Return in this EXACT format (use topic keys EXACTLY as shown above - EVERY topic listed must be in the response):
{{
  "subtopics": {{
    "topic-123": {{
      "questions": [
        {{
          "question": "Question text?",
          "options": ["Option A", "Option B", "Option C", "Option D"],
          "correctAnswer": 0,
          "explanation": "Why this answer is correct",
          "sourceText": "The complete sentence or paragraph from the study material with context.",
          "sourcePage": null
        }},
        ... (10-30+ questions for topic "topic-123")
      ]
    }},
    "topic-456": {{
      "questions": [
        {{question object}},
        ... (10-30+ questions for topic "topic-456")
      ]
    }},
    ... (MUST include ALL topic keys from the list above)
  }}
}}

REMINDER: The response MUST include questions for ALL {len(next_batch)} topics listed above. Double-check that every topic key appears in your JSON response."""

    # Make API call
    logger.info(f"📊 Prompt length: {len(batch_prompt):,} characters")

    try:
        if use_claude:
            batch_response = anthropic_client.messages.create(
                model="claude-3-5-haiku-20241022",
                max_tokens=8192,  # Maximum for Haiku (2 topics × ~20 questions each fits in 8k)
                temperature=0.7,
                messages=[{"role": "user", "content": batch_prompt}]
            )
            batch_text = batch_response.content[0].text
            logger.info(f"📊 AI response stats: stop_reason={batch_response.stop_reason}, input_tokens={batch_response.usage.input_tokens}, output_tokens={batch_response.usage.output_tokens}")
        else:
            batch_response = deepseek_client.chat.completions.create(
                model="deepseek-chat",
                max_tokens=64000,  # DeepSeek supports larger output
                temperature=0.7,
                messages=[{"role": "user", "content": batch_prompt}]
            )
            batch_text = batch_response.choices[0].message.content
    except Exception as api_error:
        logger.error(f"❌ AI API call failed: {type(api_error).__name__}: {str(api_error)}")
        raise HTTPException(status_code=500, detail=f"AI API error: {str(api_error)}")

    # Parse response
    logger.info(f"📨 Received AI response, length: {len(batch_text)} characters")

    try:
        start_idx = batch_text.find('{')
        end_idx = batch_text.rfind('}') + 1

        if start_idx == -1 or end_idx == 0:
            logger.error(f"❌ No JSON found in AI response")
            logger.error(f"❌ Full AI response (first 2000 chars):")
            logger.error(f"{batch_text[:2000]}")
            raise HTTPException(status_code=500, detail=f"AI returned non-JSON response. First 500 chars: {batch_text[:500]}")

        json_str = batch_text[start_idx:end_idx]

        # Log the JSON string for debugging
        logger.debug(f"📝 Attempting to parse JSON (length: {len(json_str)} chars)")
        logger.debug(f"📝 First 500 chars of JSON: {json_str[:500]}")

        batch_json = json.loads(json_str)
        subtopics_questions = batch_json.get("subtopics", {})

        logger.info(f"✅ Parsed {len(subtopics_questions)} topics from AI response")
        for key in subtopics_questions.keys():
            q_count = len(subtopics_questions[key].get("questions", []))
            logger.info(f"  - Topic {key}: {q_count} questions")

        # Validate that ALL topics in this batch got questions
        missing_topics = [key for key in subtopic_map.keys() if key not in subtopics_questions or not subtopics_questions[key].get("questions")]
        if missing_topics:
            logger.warning(f"⚠️ AI did NOT generate questions for {len(missing_topics)} topics: {missing_topics}")
            logger.warning(f"⚠️ Topics requested: {list(subtopic_map.keys())}")
        else:
            logger.info(f"✅ ALL {len(subtopic_map)} topics have questions")

    except json.JSONDecodeError as e:
        logger.error(f"❌ JSON parsing failed: {e}")
        logger.error(f"❌ Error at line {e.lineno}, column {e.colno}")
        logger.error(f"❌ Full JSON response (first 3000 chars):")
        logger.error(f"{json_str[:3000]}")
        # Show context around the error
        if e.pos and e.pos < len(json_str):
            start = max(0, e.pos - 200)
            end = min(len(json_str), e.pos + 200)
            logger.error(f"❌ Context around error position {e.pos}:")
            logger.error(f"{json_str[start:end]}")
        raise HTTPException(status_code=500, detail=f"Failed to parse AI JSON: {str(e)}")
    except (KeyError, ValueError) as e:
        logger.error(f"❌ Failed to parse questions: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to parse AI response: {str(e)}")

    # Save questions to database
    total_questions_generated = 0

    for key, topic in subtopic_map.items():
        questions_data = subtopics_questions.get(key, {}).get("questions", [])

        if not questions_data:
            logger.warning(f"⚠️ No questions generated for subtopic '{key}' ('{topic.title}')")
            continue

        logger.info(f"✅ Saving {len(questions_data)} questions for subtopic '{topic.title}'")

        for q_idx, q_data in enumerate(questions_data):
            # Ensure options is a list
            options = q_data.get("options", ["A", "B", "C", "D"])
            if isinstance(options, str):
                try:
                    options = json.loads(options)
                except json.JSONDecodeError:
                    options = ["Option A", "Option B", "Option C", "Option D"]

            question = Question(
                topic_id=topic.id,
                question=q_data.get("question", f"Question {q_idx+1}"),
                options=options,
                correct_answer=q_data.get("correctAnswer", 0),
                explanation=q_data.get("explanation", ""),
                source_text=q_data.get("sourceText"),
                source_page=q_data.get("sourcePage"),
                order_index=q_idx
            )
            db.add(question)
            total_questions_generated += 1

    # Commit to database
    try:
        db.commit()
        logger.info(f"✅ Successfully saved {total_questions_generated} questions to database")
    except Exception as e:
        db.rollback()
        logger.error(f"❌ Failed to save questions: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to save questions: {str(e)}")

    # Count remaining subtopics without questions
    remaining_count = len(subtopics_without_questions) - len(next_batch)

    logger.info(f"✅ Generated {total_questions_generated} questions for {len(next_batch)} subtopics")
    logger.info(f"📊 Remaining subtopics without questions: {remaining_count}")

    return {
        "message": f"Successfully generated questions for {len(next_batch)} subtopics",
        "generated": len(next_batch),
        "totalQuestions": total_questions_generated,
        "remaining": remaining_count,
        "hasMore": remaining_count > 0
    }


@router.delete("/{session_id}")
async def delete_study_session(
    session_id: str,  # UUID as string
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db),
):
    """
    Delete a study session and all its associated data.

    This will cascade delete all topics and questions associated with the session.
    """
    # Validate UUID format
    try:
        uuid_obj = uuid.UUID(session_id)
    except ValueError:
        logger.error(f"❌ Invalid session ID format: {session_id} (expected UUID)")
        raise HTTPException(
            status_code=400,
            detail=f"Invalid session ID format. Expected UUID, got: '{session_id}'."
        )

    session = db.query(StudySession).filter(
        StudySession.id == uuid_obj,
        StudySession.user_id == current_user.id
    ).first()

    if not session:
        raise HTTPException(status_code=404, detail="Study session not found")

    db.delete(session)
    db.commit()

    return {"message": "Study session deleted successfully", "id": session_id}


@router.patch("/{session_id}/archive")
async def archive_study_session(
    session_id: str,  # UUID as string
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db),
):
    """
    Archive a study session.

    Changes the session status to 'archived'.
    """
    # Validate UUID format
    try:
        uuid_obj = uuid.UUID(session_id)
    except ValueError:
        logger.error(f"❌ Invalid session ID format: {session_id} (expected UUID)")
        raise HTTPException(
            status_code=400,
            detail=f"Invalid session ID format. Expected UUID, got: '{session_id}'."
        )

    session = db.query(StudySession).filter(
        StudySession.id == uuid_obj,
        StudySession.user_id == current_user.id
    ).first()

    if not session:
        raise HTTPException(status_code=404, detail="Study session not found")

    session.status = "archived"
    db.commit()
    db.refresh(session)

    return {"message": "Study session archived successfully", "id": session_id}


class UpdateTopicProgressRequest(BaseModel):
    """Request schema for updating topic progress."""
    score: int = Field(..., ge=0, le=100)  # Score as percentage 0-100
    current_question_index: int = Field(..., ge=0)
    completed: bool = False


class UpdateUserXPRequest(BaseModel):
    """Request schema for updating user XP."""
    xp_to_add: int = Field(..., ge=0, le=1000)  # XP to add (capped at 1000 per call)


@router.patch("/{session_id}/topics/{topic_id}/progress")
@limiter.limit("60/minute")
async def update_topic_progress(
    request: Request,
    session_id: str,
    topic_id: int,  # Database topic ID
    data: UpdateTopicProgressRequest,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db),
):
    """
    Update topic progress (score, current question index, completion status).

    This endpoint is called after each answer to persist user progress.

    Rate Limits:
        - 60 requests per minute per user
    """
    # Validate UUID format
    try:
        uuid_obj = uuid.UUID(session_id)
    except ValueError:
        logger.error(f"❌ Invalid session ID format: {session_id} (expected UUID)")
        raise HTTPException(
            status_code=400,
            detail=f"Invalid session ID format. Expected UUID, got: '{session_id}'."
        )

    # Verify session belongs to user
    session = db.query(StudySession).filter(
        StudySession.id == uuid_obj,
        StudySession.user_id == current_user.id
    ).first()

    if not session:
        raise HTTPException(status_code=404, detail="Study session not found")

    # Find the topic
    topic = db.query(Topic).filter(
        Topic.id == topic_id,
        Topic.study_session_id == uuid_obj
    ).first()

    if not topic:
        raise HTTPException(status_code=404, detail="Topic not found")

    # Update topic progress
    topic.score = data.score
    topic.current_question_index = data.current_question_index
    topic.completed = data.completed

    # Update session progress (percentage of completed subtopics)
    all_topics = db.query(Topic).filter(
        Topic.study_session_id == uuid_obj,
        Topic.is_category == False  # Only count leaf topics
    ).all()

    completed_topics = sum(1 for t in all_topics if t.completed)
    total_topics = len(all_topics)
    session.progress = int((completed_topics / total_topics * 100) if total_topics > 0 else 0)

    db.commit()
    db.refresh(topic)

    return {
        "message": "Topic progress updated successfully",
        "topic_id": topic_id,
        "score": topic.score,
        "current_question_index": topic.current_question_index,
        "completed": topic.completed,
        "session_progress": session.progress
    }


@router.patch("/user/xp")
@limiter.limit("100/minute")
async def update_user_xp(
    request: Request,
    data: UpdateUserXPRequest,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db),
):
    """
    Add XP to the current user's total.

    This endpoint is called after each correct answer.

    Rate Limits:
        - 100 requests per minute per user
    """
    # Update user XP
    current_user.xp += data.xp_to_add
    current_user.updated_at = datetime.utcnow()

    # Calculate level (simple formula: level = floor(xp / 100) + 1)
    # Every 100 XP = 1 level
    current_user.level = (current_user.xp // 100) + 1

    db.commit()
    db.refresh(current_user)

    return {
        "message": "XP updated successfully",
        "xp": current_user.xp,
        "xp_added": data.xp_to_add,
        "level": current_user.level
    }
